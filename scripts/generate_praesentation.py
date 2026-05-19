#!/usr/bin/env python3
"""
generate_praesentation.py  –  Trainer-Präsentation als PPTX
Design: FKB Campus Masterdeck (editorial, Variante B)
Verwendung: python3 scripts/generate_praesentation.py M01
            python3 scripts/generate_praesentation.py all
Ausgabe:    public/downloads/praesentation/MXX.pptx
"""

import sys, os, re, textwrap
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree
    import yaml
except ImportError as e:
    print(f"Fehlende Abhängigkeit: {e}\npip3 install python-pptx pyyaml lxml")
    sys.exit(1)

# ── Farbsystem (FKB Campus Masterdeck, oklch → hex) ───────────────────────────
BG       = RGBColor(0xF8, 0xF8, 0xFA)   # off-white
BG2      = RGBColor(0xF0, 0xF0, 0xF5)   # sekundär
INK      = RGBColor(0x19, 0x1D, 0x2E)   # fast schwarz, leicht blau
INK2     = RGBColor(0x4D, 0x52, 0x70)   # body text
INK3     = RGBColor(0x7B, 0x7F, 0xA0)   # meta, labels
LINE     = RGBColor(0xE0, 0xE0, 0xE8)   # standard trennlinie
LINE2    = RGBColor(0xC8, 0xC8, 0xD5)   # stärkere trennlinie
PRIMARY  = RGBColor(0x1F, 0x2B, 0x56)   # genossenschafts-tiefblau
PRI_INK  = RGBColor(0xF8, 0xF8, 0xFA)   # text auf primary
ACCENT   = RGBColor(0xD9, 0xBF, 0x7A)   # warmer sand
ACC_INK  = RGBColor(0x3A, 0x2A, 0x0A)   # text auf accent
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# ── Typografie ─────────────────────────────────────────────────────────────────
SERIF = "Georgia"          # Source Serif 4 → Georgia (universell verfügbar)
SANS  = "Arial"            # Inter Tight → Arial
MONO  = "Courier New"      # JetBrains Mono → Courier New

# ── Foliengröße: 1920×1080 (16:9, widescreen) ─────────────────────────────────
# 1920px / 144 px/inch = 13.333"; 1080px / 144 = 7.5"
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
PX_PER_IN  = 144  # skalierungsfaktor

def px(n): return Inches(n / PX_PER_IN)

# ── Pfade ──────────────────────────────────────────────────────────────────────
SCRIPT_DIR   = Path(__file__).parent
AKADEMIE_DIR = SCRIPT_DIR.parent
MODULES_DIR  = AKADEMIE_DIR / "content" / "modules"
OUTPUT_DIR   = AKADEMIE_DIR / "public" / "downloads" / "praesentation"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


# ══════════════════════════════════════════════════════════════════════════════
# MODUL PARSEN
# ══════════════════════════════════════════════════════════════════════════════

def parse_module(module_id: str) -> dict:
    path = MODULES_DIR / f"{module_id.upper()}.md"
    if not path.exists():
        print(f"Modul nicht gefunden: {path}"); sys.exit(1)

    text = path.read_text(encoding="utf-8")
    fm_m = re.match(r"^---\s*\n(.*?)\n---\s*\n", text, re.DOTALL)
    if not fm_m:
        print("Kein Frontmatter."); sys.exit(1)

    fm   = yaml.safe_load(fm_m.group(1))
    body = text[fm_m.end():]

    def extract_block(key):
        m = re.search(
            rf"<!-- CONTENT_{key}_START -->(.*?)<!-- CONTENT_{key}_END -->",
            body, re.DOTALL)
        return (m.group(1).strip(), body[:m.start()] + body[m.end():]) if m else ("", body)

    trainer_content, body = extract_block("TRAINER")
    _,               body = extract_block("THEORIE")
    body = re.sub(r"<!--.*?-->", "", body, flags=re.DOTALL)
    body = re.sub(r"^>.*?Stand:.*$", "", body, flags=re.MULTILINE)
    body = re.sub(r"\n{3,}", "\n\n", body)

    lernziele = fm.get("lernziele", [])
    lz_texts  = [lz["text"] if isinstance(lz, dict) else str(lz) for lz in lernziele]

    return {
        "id":            str(fm.get("id", module_id)),
        "title":         str(fm.get("title", "")),
        "subtitle":      str(fm.get("subtitle", "")),
        "kompetenzfeld": str(fm.get("kompetenzfeld", "")),
        "stufe":         str(fm.get("stufe", "")),
        "dauer":         str(fm.get("dauer", "")),
        "version":       str(fm.get("version", "")),
        "lernziele":     lz_texts,
        "body":          body.strip(),
        "trainer":       trainer_content,
    }


# ══════════════════════════════════════════════════════════════════════════════
# MARKDOWN-PARSER
# ══════════════════════════════════════════════════════════════════════════════

def clean_md(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*",     r"\1", text)
    text = re.sub(r"`(.+?)`",       r"\1", text)
    text = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", text)
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    return text.strip()

def extract_bullets(text: str, max_items=6) -> list:
    bullets = []
    for line in text.splitlines():
        line = line.strip()
        if re.match(r"^[-*•]\s+", line):
            b = re.sub(r"^[-*•]\s+", "", line)
            b = clean_md(b)
            if b: bullets.append(b)
    if not bullets:
        for para in re.split(r"\n{2,}", text):
            para = clean_md(para).strip()
            if para and not para.startswith("|") and len(para) > 20:
                bullets.append(textwrap.shorten(para, 200, placeholder="…"))
    return bullets[:max_items]

def parse_md_table(text: str):
    rows = []
    for line in text.splitlines():
        line = line.strip()
        if not line.startswith("|"): continue
        if re.match(r"^\|[-| :]+\|$", line): continue
        cells = [clean_md(c.strip()) for c in line.strip("|").split("|")]
        rows.append(cells)
    return rows if len(rows) >= 2 else None

def split_sections(body: str, level: int) -> list:
    pattern = rf"^{'#'*level}\s+(.+)$"
    parts = re.split(pattern, body, flags=re.MULTILINE)
    result = []
    for i in range(1, len(parts)-1, 2):
        result.append({"title": parts[i].strip(), "body": parts[i+1].strip()})
    return result

def find_praxisfall(body: str) -> tuple:
    """Gibt (name, situation_text) zurück, falls Praxisfall gefunden."""
    m = re.search(r"###\s+(?:Praxisfall|Der Fall|Fallstudie)[^\n]*\n(.*?)(?=###|\Z)",
                  body, re.DOTALL | re.IGNORECASE)
    if not m: return None, None
    content = m.group(1).strip()
    # Firmenname aus erstem Bold-Element
    nm = re.search(r"\*\*(.+?)\*\*", content)
    name = nm.group(1) if nm else "Praxisfall"
    situation = clean_md(content[:500])
    return name, situation

def find_reflexionsfragen(body: str) -> list:
    m = re.search(r"###\s+(?:Reflexionsfragen|Fragen zur Reflexion)[^\n]*\n(.*?)(?=###|\Z)",
                  body, re.DOTALL | re.IGNORECASE)
    if not m: return []
    questions = []
    for line in m.group(1).splitlines():
        line = line.strip()
        qm = re.match(r"^\d+\.\s+(.+)$", line)
        if qm: questions.append(qm.group(1).strip())
    return questions[:5]


# ══════════════════════════════════════════════════════════════════════════════
# FOLIE-BASISWERKZEUGE
# ══════════════════════════════════════════════════════════════════════════════

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs

def blank(prs) -> object:
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_slide(slide, color: RGBColor):
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = color

def shape_rect(slide, l, t, w, h, color: RGBColor):
    """Rechteck ohne Linie hinzufügen. Koordinaten in Inch."""
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def shape_rect_px(slide, l, t, w, h, color: RGBColor):
    """Rechteck mit Pixel-Koordinaten."""
    return shape_rect(slide, l/PX_PER_IN, t/PX_PER_IN, w/PX_PER_IN, h/PX_PER_IN, color)

def rule_px(slide, l, t, w, color: RGBColor):
    """Horizontale Linie (1px Höhe)."""
    shape_rect_px(slide, l, t, w, 1, color)

def logo_runner(slide, bg_dark=False, section_label="", page_no=""):
    """Runner-Top: Logo links, Abschnittsbezeichnung rechts."""
    # Farbe
    c = PRI_INK if bg_dark else INK3
    l_color = PRI_INK if bg_dark else INK

    # Logo-Text (ohne SVG in PPTX – stattdessen Text-Simulation)
    _tbox(slide, "FKB  Campus", px_l=96, px_t=48, px_w=400, px_h=44,
          font=SANS, size=16, bold=True, color=l_color, align=PP_ALIGN.LEFT)

    if section_label:
        _tbox(slide, section_label.upper(), px_l=0, px_t=48, px_w=1824, px_h=44,
              font=MONO, size=13, color=c, align=PP_ALIGN.RIGHT,
              letter_spacing=True)

def runner_bot(slide, bg_dark=False, edition="", page_no=""):
    c = RGBColor(0x60, 0x60, 0x70) if bg_dark else INK3
    if edition:
        _tbox(slide, edition, px_l=96, px_t=1000, px_w=800, px_h=36,
              font=MONO, size=13, color=c, align=PP_ALIGN.LEFT)
    if page_no:
        _tbox(slide, page_no, px_l=0, px_t=1000, px_w=1824, px_h=36,
              font=MONO, size=13, color=c, align=PP_ALIGN.RIGHT)


def _tbox(slide, text: str, px_l=0, px_t=0, px_w=200, px_h=60,
          font=SANS, size=20, bold=False, italic=False,
          color=INK, align=PP_ALIGN.LEFT,
          letter_spacing=False, wrap=True) -> None:
    """Text-Box mit Pixel-Koordinaten."""
    if not text: return
    tx = slide.shapes.add_textbox(
        px(px_l), px(px_t), px(px_w), px(px_h))
    tf = tx.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name   = font
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if letter_spacing:
        # Laufweite über XML
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", "100")  # 100/1000 em ≈ 0.1em

def _tbox_multiline(slide, lines: list, px_l, px_t, px_w, px_h,
                    font=SANS, size=18, bold=False, italic=False,
                    color=INK, line_spacing=1.3, bullet_char=""):
    """Mehrere Zeilen als separate Paragraphen in einer Textbox."""
    if not lines: return
    tx = slide.shapes.add_textbox(px(px_l), px(px_t), px(px_w), px(px_h))
    tf = tx.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(size * (line_spacing - 1) * 0.5)
        p.space_after  = Pt(size * (line_spacing - 1) * 0.5)
        run = p.add_run()
        if bullet_char:
            run.text = f"{bullet_char}  {line}"
        else:
            run.text = str(line)
        run.font.name   = font
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════════════════
# FOLIEN-TYPEN (nach FKB Campus Masterdeck)
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs, m, page_total=""):
    """01 — COVER: großer Titel unten links, Meta-Grid."""
    slide = blank(prs)
    fill_slide(slide, BG)

    # Edition-Label oben
    _tbox(slide, "FKB  Campus", px_l=96, px_t=48, px_w=500, px_h=44,
          font=SANS, size=16, bold=True, color=INK)
    _tbox(slide, f"{m['id']}  ·  {m['version']}",
          px_l=0, px_t=48, px_w=1824, px_h=44,
          font=MONO, size=14, color=INK3, align=PP_ALIGN.RIGHT)

    # Großer Titel
    title = m["title"]
    _tbox(slide, title, px_l=96, px_t=480, px_w=1728, px_h=340,
          font=SERIF, size=96, bold=False, color=INK, align=PP_ALIGN.LEFT)

    # Meta-Trennlinie
    rule_px(slide, 96, 840, 1728, INK)

    # Meta-Grid: 3 Spalten
    # Spalte 1: Untertitel/Lead
    sub = m.get("subtitle", "") or m.get("kompetenzfeld", "")
    _tbox(slide, sub, px_l=96, px_t=860, px_w=820, px_h=140,
          font=SERIF, size=22, italic=True, color=INK2)

    # Spalte 2: Zielstufe
    _tbox(slide, "ZIELSTUFE", px_l=950, px_t=860, px_w=360, px_h=28,
          font=MONO, size=12, color=INK3, letter_spacing=True)
    _tbox(slide, m["stufe"], px_l=950, px_t=892, px_w=360, px_h=80,
          font=SERIF, size=22, color=INK)

    # Spalte 3: Dauer
    _tbox(slide, "DAUER", px_l=1380, px_t=860, px_w=400, px_h=28,
          font=MONO, size=12, color=INK3, letter_spacing=True)
    _tbox(slide, m["dauer"], px_l=1380, px_t=892, px_w=400, px_h=80,
          font=SERIF, size=22, color=INK)


def slide_02_agenda(prs, m, page_no=2, page_total=""):
    """02 — AGENDA: Lernziele als strukturierte Liste."""
    slide = blank(prs)
    fill_slide(slide, BG)
    logo_runner(slide, section_label="§ 00 — Lernziele")
    runner_bot(slide, edition=m["version"],
               page_no=f"{page_no:02d} / {page_total}")

    # Linke Spalte: Titel
    _tbox(slide, "Inhalt.", px_l=96, px_t=200, px_w=580, px_h=200,
          font=SERIF, size=88, italic=True, color=PRIMARY)

    # Trennlinie Top
    rule_px(slide, 96, 210, 1728, INK)

    # Rechte Spalte: Nummerierte Lernziele
    lz = m["lernziele"]
    romans = ["I.", "II.", "III.", "IV.", "V.", "VI.", "VII.", "VIII."]
    top = 210
    row_h = min(120, (840 - top) // max(len(lz), 1))

    for i, lz_text in enumerate(lz[:7]):
        y = top + i * row_h
        rule_px(slide, 720, y, 1104, LINE)
        # Römische Zahl
        _tbox(slide, romans[i] if i < len(romans) else str(i+1),
              px_l=720, px_t=y+8, px_w=90, px_h=row_h-8,
              font=SERIF, size=40, italic=True, color=PRIMARY)
        # Lernzieltext
        text = textwrap.shorten(lz_text, 160, placeholder="…")
        _tbox(slide, text, px_l=830, px_t=y+12, px_w=980, px_h=row_h-8,
              font=SERIF, size=22, color=INK)
        # Zeit
        _tbox(slide, "", px_l=1750, px_t=y+16, px_w=60, px_h=36,
              font=MONO, size=14, color=INK3, align=PP_ALIGN.RIGHT)


def slide_03_section(prs, roman_num: str, title: str, subtitle: str,
                     section_label: str, page_no: int, page_total: str):
    """03 — SECTION DIVIDER: dunkler Hintergrund, riesige Zahl links."""
    slide = blank(prs)
    fill_slide(slide, INK)
    logo_runner(slide, bg_dark=True, section_label=section_label)
    runner_bot(slide, bg_dark=True, page_no=f"{page_no:02d} / {page_total}")

    # Riesige Zahl links
    _tbox(slide, roman_num, px_l=96, px_t=130, px_w=900, px_h=780,
          font=SERIF, size=400, italic=True, color=ACCENT)

    # Titel und Text rechts unten
    _tbox(slide, section_label.upper(), px_l=1100, px_t=580, px_w=724, px_h=36,
          font=MONO, size=14, color=INK3, align=PP_ALIGN.RIGHT)
    _tbox(slide, title, px_l=1100, px_t=630, px_w=724, px_h=240,
          font=SERIF, size=76, color=PRI_INK, align=PP_ALIGN.RIGHT)
    if subtitle:
        sub = textwrap.shorten(subtitle, 200, placeholder="…")
        _tbox(slide, sub, px_l=1100, px_t=870, px_w=724, px_h=100,
              font=SERIF, size=20, italic=True, color=INK3, align=PP_ALIGN.RIGHT)


def slide_04_headline(prs, title: str, body_left: str, body_right: str,
                      section_label: str, page_no: int, page_total: str):
    """04 — HEADLINE + STANDFIRST: große Überschrift, zweispaltig darunter."""
    slide = blank(prs)
    fill_slide(slide, BG)
    logo_runner(slide, section_label=section_label)
    runner_bot(slide, page_no=f"{page_no:02d} / {page_total}")

    # Eyebrow
    _tbox(slide, "Kernaussage".upper(), px_l=96, px_t=140, px_w=600, px_h=32,
          font=MONO, size=13, color=INK3, letter_spacing=True)

    # Headline
    _tbox(slide, title, px_l=96, px_t=180, px_w=1728, px_h=360,
          font=SERIF, size=88, color=INK)

    # Trennlinie
    rule_px(slide, 96, 560, 1728, INK)

    # Zweispaltig darunter
    if body_left:
        text = textwrap.shorten(body_left, 350, placeholder="…")
        _tbox(slide, text, px_l=96, px_t=580, px_w=820, px_h=340,
              font=SERIF, size=22, color=INK, italic=False)

    if body_right:
        text = textwrap.shorten(body_right, 300, placeholder="…")
        _tbox(slide, "Mitnahme".upper(), px_l=980, px_t=580, px_w=740, px_h=32,
              font=MONO, size=12, color=PRIMARY, letter_spacing=True)
        _tbox(slide, text, px_l=980, px_t=620, px_w=740, px_h=300,
              font=SERIF, size=22, italic=True, color=INK)


def slide_05_two_col(prs, title: str, left_items: list, right_items: list,
                     section_label: str, page_no: int, page_total: str):
    """05 — TWO-COLUMN BODY: Titel oben, zwei Spalten mit Bullet-Items."""
    slide = blank(prs)
    fill_slide(slide, BG)
    logo_runner(slide, section_label=section_label)
    runner_bot(slide, page_no=f"{page_no:02d} / {page_total}")

    # Titel über beide Spalten
    _tbox(slide, title, px_l=96, px_t=140, px_w=1728, px_h=200,
          font=SERIF, size=64, color=INK)
    rule_px(slide, 96, 340, 1728, INK)

    # Linke Spalte
    left_label = left_items[0] if left_items else ""
    left_body  = left_items[1:] if len(left_items) > 1 else left_items

    if left_label:
        _tbox(slide, left_label.upper(), px_l=96, px_t=360, px_w=800, px_h=32,
              font=MONO, size=12, color=INK3, letter_spacing=True)
        rule_px(slide, 96, 396, 800, INK)

    for i, item in enumerate(left_body[:5]):
        y = 400 + i * 100
        _tbox(slide, f"0{i+1}", px_l=96, px_t=y, px_w=64, px_h=40,
              font=MONO, size=16, color=PRIMARY)
        text = textwrap.shorten(item, 160, placeholder="…")
        _tbox(slide, text, px_l=176, px_t=y, px_w=700, px_h=88,
              font=SERIF, size=20, color=INK)
        rule_px(slide, 96, y+90, 800, LINE)

    # Rechte Spalte
    right_label = right_items[0] if right_items else ""
    right_body  = right_items[1:] if len(right_items) > 1 else right_items

    if right_label:
        _tbox(slide, right_label.upper(), px_l=1024, px_t=360, px_w=800, px_h=32,
              font=MONO, size=12, color=INK3, letter_spacing=True)
        rule_px(slide, 1024, 396, 800, INK)

    for i, item in enumerate(right_body[:5]):
        y = 400 + i * 100
        _tbox(slide, f"0{i+1}", px_l=1024, px_t=y, px_w=64, px_h=40,
              font=MONO, size=16, color=PRIMARY)
        text = textwrap.shorten(item, 160, placeholder="…")
        _tbox(slide, text, px_l=1104, px_t=y, px_w=700, px_h=88,
              font=SERIF, size=20, color=INK)
        rule_px(slide, 1024, y+90, 800, LINE)


def slide_06_pullquote(prs, quote: str, author: str, role: str,
                       section_label: str, page_no: int, page_total: str):
    """06 — PULL QUOTE: großes Zitat auf primary-Hintergrund."""
    slide = blank(prs)
    fill_slide(slide, PRIMARY)
    logo_runner(slide, bg_dark=True, section_label=section_label)
    runner_bot(slide, bg_dark=True, page_no=f"{page_no:02d} / {page_total}")

    # Anführungszeichen
    _tbox(slide, "„", px_l=96, px_t=160, px_w=400, px_h=300,
          font=SERIF, size=240, italic=True, color=ACCENT)

    # Zitat
    q = textwrap.shorten(quote, 280, placeholder="…")
    _tbox(slide, q, px_l=96, px_t=400, px_w=1728, px_h=380,
          font=SERIF, size=64, italic=True, color=PRI_INK)

    # Trennlinie Attribution
    rule_px(slide, 96, 800, 720, RGBColor(0x40, 0x48, 0x70))

    # Initialen-Avatar
    initials = "".join(w[0] for w in author.split()[:2]).upper() if author else "??"
    shape_rect_px(slide, 96, 824, 72, 72, ACCENT)
    _tbox(slide, initials, px_l=96, px_t=832, px_w=72, px_h=56,
          font=MONO, size=20, bold=True, color=ACC_INK, align=PP_ALIGN.CENTER)

    _tbox(slide, author, px_l=192, px_t=824, px_w=624, px_h=40,
          font=SERIF, size=24, bold=True, color=PRI_INK)
    _tbox(slide, role, px_l=192, px_t=866, px_w=624, px_h=36,
          font=MONO, size=13, color=INK3, letter_spacing=True)


def slide_07_stats(prs, title: str, stats: list, context: str,
                   section_label: str, page_no: int, page_total: str):
    """07 — STATS GRID: Headline + 2×2 Kennzahlen."""
    slide = blank(prs)
    fill_slide(slide, BG)
    logo_runner(slide, section_label=section_label)
    runner_bot(slide, page_no=f"{page_no:02d} / {page_total}")

    # Headline links
    _tbox(slide, title, px_l=96, px_t=140, px_w=900, px_h=200,
          font=SERIF, size=72, color=INK)
    if context:
        _tbox(slide, textwrap.shorten(context, 200, placeholder="…"),
              px_l=1000, px_t=240, px_w=824, px_h=100,
              font=SERIF, size=20, color=INK2)

    rule_px(slide, 96, 350, 1728, INK)

    # 2×2 Kacheln
    positions = [(96, 370), (960, 370), (96, 700), (960, 700)]
    for i, (stat_val, stat_label) in enumerate(stats[:4]):
        x, y = positions[i]
        _tbox(slide, str(stat_val), px_l=x, px_t=y, px_w=840, px_h=220,
              font=SERIF, size=144, italic=False, color=INK)
        _tbox(slide, str(stat_label), px_l=x, px_t=y+200, px_w=840, px_h=60,
              font=MONO, size=14, color=INK3, letter_spacing=True)


def slide_08_table(prs, title: str, rows: list, context: str,
                   section_label: str, page_no: int, page_total: str):
    """08 — DATA TABLE: Titel links, Tabelle rechts."""
    from pptx.oxml import parse_xml
    slide = blank(prs)
    fill_slide(slide, BG)
    logo_runner(slide, section_label=section_label)
    runner_bot(slide, page_no=f"{page_no:02d} / {page_total}")

    # Linke Spalte
    _tbox(slide, title, px_l=96, px_t=160, px_w=780, px_h=320,
          font=SERIF, size=60, color=INK)
    if context:
        _tbox(slide, textwrap.shorten(context, 200, placeholder="…"),
              px_l=96, px_t=500, px_w=780, px_h=200,
              font=SERIF, size=20, color=INK2)

    # Rechts: Tabelle
    if not rows: return
    nr  = len(rows)
    nc  = max(len(r) for r in rows)
    tw, th = Inches(1024/PX_PER_IN), Inches(min(600, nr * 70) / PX_PER_IN)
    tbl = slide.shapes.add_table(nr, nc, px(960), px(160), tw, th).table

    for ri, row in enumerate(rows):
        for ci in range(nc):
            cell = tbl.cell(ri, ci)
            text = row[ci] if ci < len(row) else ""
            tf   = cell.text_frame
            tf.word_wrap = True
            p    = tf.paragraphs[0]
            run  = p.add_run()
            run.text = clean_md(text)
            run.font.name  = MONO
            run.font.size  = Pt(14 if ri == 0 else 16)
            run.font.bold  = (ri == 0)

            # Farben per XML
            if ri == 0:
                run.font.color.rgb = BG
                _set_cell_bg(cell, INK)
            elif ri % 2 == 0:
                run.font.color.rgb = INK
                _set_cell_bg(cell, BG2)
            else:
                run.font.color.rgb = INK
                _set_cell_bg(cell, BG)


def _set_cell_bg(cell, color: RGBColor):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    hex6 = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"
    fill = etree.fromstring(
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{hex6}"/></a:solidFill>')
    tcPr.append(fill)


def slide_09_speakers(prs, title: str, people: list,
                      section_label: str, page_no: int, page_total: str):
    """09 — SPEAKERS / AUTOREN: zwei Personen nebeneinander."""
    slide = blank(prs)
    fill_slide(slide, BG)
    logo_runner(slide, section_label=section_label)
    runner_bot(slide, page_no=f"{page_no:02d} / {page_total}")

    _tbox(slide, title, px_l=96, px_t=140, px_w=1728, px_h=220,
          font=SERIF, size=88, color=INK)
    rule_px(slide, 96, 360, 1728, INK)

    # Personenkarten
    col_x = [96, 980]
    for i, p in enumerate(people[:2]):
        x = col_x[i]
        name = p.get("name", "")
        role = p.get("role", "")
        bio  = p.get("bio", "")

        # Portrait-Platzhalter (schraffiertes Rechteck)
        shape_rect_px(slide, x, 380, 200, 240, BG2)
        # Diagonale Linien simulieren (einfaches Rechteck mit Outline)
        s = slide.shapes.add_shape(1, px(x), px(380), px(200), px(240))
        s.fill.solid(); s.fill.fore_color.rgb = BG2
        s.line.color.rgb = LINE2

        _tbox(slide, name, px_l=x+220, px_t=380, px_w=640, px_h=100,
              font=SERIF, size=34, bold=False, color=INK)
        _tbox(slide, role.upper(), px_l=x+220, px_t=490, px_w=640, px_h=32,
              font=MONO, size=13, color=PRIMARY, letter_spacing=True)
        _tbox(slide, textwrap.shorten(bio, 200, placeholder="…"),
              px_l=x+220, px_t=530, px_w=640, px_h=150,
              font=SERIF, size=18, color=INK2)


def slide_10_image_led(prs, title: str, situation: str, label: str,
                       section_label: str, page_no: int, page_total: str):
    """10 — IMAGE-LED: linkes Bilddrittel (Platzhalter), rechts Text."""
    slide = blank(prs)
    fill_slide(slide, BG)
    logo_runner(slide, section_label=section_label)
    runner_bot(slide, page_no=f"{page_no:02d} / {page_total}")

    # Bild-Platzhalter (schraffierter Bereich in Accent-Tönen)
    shape_rect_px(slide, 0, 0, 740, 1080, RGBColor(0xCA, 0xB0, 0x68))

    # Schraffur-Simulation: mehrere dünne Linien
    for yi in range(0, 1080, 24):
        shape_rect_px(slide, 0, yi, 740, 1, RGBColor(0xBE, 0xA2, 0x55))

    # Bildunterschrift
    _tbox(slide, label.upper(), px_l=24, px_t=1020, px_w=700, px_h=40,
          font=MONO, size=12, color=ACC_INK, letter_spacing=True)

    # Rechte Spalte
    _tbox(slide, "Praxisfall".upper(), px_l=796, px_t=200, px_w=1028, px_h=32,
          font=MONO, size=13, color=INK3, letter_spacing=True)
    _tbox(slide, title, px_l=796, px_t=240, px_w=1028, px_h=320,
          font=SERIF, size=64, color=INK)
    if situation:
        sit = textwrap.shorten(situation, 400, placeholder="…")
        _tbox(slide, sit, px_l=796, px_t=580, px_w=1028, px_h=300,
              font=SERIF, size=22, color=INK2)


def slide_11_sources(prs, sources: list, page_no: int, page_total: str):
    """11 — SOURCES / APPARAT: dunkler Hintergrund, 2-spaltige Quellenliste."""
    slide = blank(prs)
    fill_slide(slide, INK)
    logo_runner(slide, bg_dark=True, section_label="Quellenapparat")
    runner_bot(slide, bg_dark=True, page_no=f"{page_no:02d} / {page_total}")

    _tbox(slide, "Quellenapparat".upper(), px_l=96, px_t=140, px_w=800, px_h=36,
          font=MONO, size=14, color=INK3, letter_spacing=True)
    _tbox(slide, "Quellen­apparat.", px_l=96, px_t=185, px_w=1728, px_h=160,
          font=SERIF, size=72, color=PRI_INK)
    rule_px(slide, 96, 360, 1728, RGBColor(0x40, 0x48, 0x70))

    mid = len(sources) // 2 + 1
    left  = sources[:mid]
    right = sources[mid:]

    for i, src in enumerate(left[:7]):
        y = 380 + i * 80
        _tbox(slide, src, px_l=96, px_t=y, px_w=800, px_h=72,
              font=MONO, size=14, color=RGBColor(0xC0, 0xC4, 0xD8))

    for i, src in enumerate(right[:7]):
        y = 380 + i * 80
        _tbox(slide, src, px_l=960, px_t=y, px_w=800, px_h=72,
              font=MONO, size=14, color=RGBColor(0xC0, 0xC4, 0xD8))


def slide_12_closing(prs, m, page_no: int, page_total: str):
    """12 — CLOSING: riesiger Titel, Kontakt-Grid unten."""
    slide = blank(prs)
    fill_slide(slide, BG)

    # Logo oben links
    _tbox(slide, "FKB  Campus", px_l=96, px_t=80, px_w=500, px_h=48,
          font=SANS, size=18, bold=True, color=INK)
    _tbox(slide, f"Ende · {m['version']}", px_l=0, px_t=86, px_w=1824, px_h=40,
          font=MONO, size=14, color=INK3, align=PP_ALIGN.RIGHT)

    # Riesiger Titel
    _tbox(slide, "Danke.\nFragen?", px_l=96, px_t=180, px_w=1728, px_h=580,
          font=SERIF, size=180, italic=False, color=INK)

    # Trennlinie
    rule_px(slide, 96, 810, 1728, INK)

    # Unteres Grid: 3 Spalten
    lead = (f"Das vollständige Modul «{m['title']}» ist auf FKB Campus verfügbar "
            f"— {len(m['lernziele'])} Lernziele, Praxisfall und Quellenapparat.")
    _tbox(slide, lead, px_l=96, px_t=830, px_w=900, px_h=160,
          font=SERIF, size=22, italic=True, color=INK)

    _tbox(slide, "Kontakt", px_l=1050, px_t=830, px_w=380, px_h=32,
          font=MONO, size=13, color=INK3, letter_spacing=True)
    _tbox(slide, "info@benedikt-zoller.de", px_l=1050, px_t=866, px_w=380, px_h=60,
          font=SERIF, size=20, color=INK)

    _tbox(slide, "Modul", px_l=1480, px_t=830, px_w=344, px_h=32,
          font=MONO, size=13, color=INK3, letter_spacing=True)
    _tbox(slide, f"{m['id']} · {m['stufe']}", px_l=1480, px_t=866, px_w=344, px_h=60,
          font=SERIF, size=20, color=INK)


# Hilfs-Folie: Bullet-Content
def slide_content(prs, title: str, bullets: list, section_label: str,
                  page_no: int, page_total: str):
    """Standard-Content-Folie: Titel + Bullet-Liste."""
    slide = blank(prs)
    fill_slide(slide, BG)
    logo_runner(slide, section_label=section_label)
    runner_bot(slide, page_no=f"{page_no:02d} / {page_total}")

    _tbox(slide, title, px_l=96, px_t=140, px_w=1728, px_h=220,
          font=SERIF, size=72, color=INK)
    rule_px(slide, 96, 360, 1728, INK)

    for i, b in enumerate(bullets[:7]):
        y = 380 + i * 100
        _tbox(slide, "–", px_l=96, px_t=y+4, px_w=48, px_h=56,
              font=SERIF, size=28, color=PRIMARY)
        text = textwrap.shorten(b, 200, placeholder="…")
        _tbox(slide, text, px_l=160, px_t=y, px_w=1560, px_h=90,
              font=SERIF, size=26, color=INK)


# ══════════════════════════════════════════════════════════════════════════════
# MODUL → PRÄSENTATION
# ══════════════════════════════════════════════════════════════════════════════

ROMAN = ["I.", "II.", "III.", "IV.", "V.", "VI.", "VII.", "VIII.", "IX.", "X."]

def build_prs(module_id: str) -> Path:
    m   = parse_module(module_id)
    prs = new_prs()

    body    = m["body"]
    h2s     = split_sections(body, 2)

    # Seiten zählen (Schätzung)
    estimated = 3 + len(h2s) * 2 + 3
    pt = str(estimated)

    page = 1

    # 01 — COVER
    slide_01_cover(prs, m, pt)
    page += 1

    # 02 — AGENDA (Lernziele)
    if m["lernziele"]:
        slide_02_agenda(prs, m, page, pt)
        page += 1

    romans_used = 0

    for h2 in h2s:
        h2_title = clean_md(h2["title"])
        h2_body  = h2["body"]

        # Sektionen ohne inhaltlichen Wert überspringen
        skip_kw = ["quellenverzeichnis", "literaturverzeichnis", "quellen"]
        if any(k in h2_title.lower() for k in skip_kw):
            # Quellenapparat extrahieren
            src_lines = [l.strip("- ").strip() for l in h2_body.splitlines()
                         if l.strip() and not l.startswith("#")]
            if src_lines:
                slide_11_sources(prs, [clean_md(s) for s in src_lines[:14]], page, pt)
                page += 1
            continue

        # 03 — SECTION DIVIDER
        roman = ROMAN[romans_used] if romans_used < len(ROMAN) else f"{romans_used+1}."
        romans_used += 1
        h3s = split_sections(h2_body, 3)
        sub = clean_md(h3s[0]["title"]) if h3s else ""
        slide_03_section(prs, roman, h2_title, sub, h2_title, page, pt)
        page += 1

        # Praxisfall erkennen
        pf_name, pf_sit = find_praxisfall(h2_body)

        for h3 in h3s:
            h3_title = clean_md(h3["title"])
            h3_body  = h3["body"]
            sec_label = f"§ {roman[:-1]} · {h2_title}"

            # Musterlösung / Trainer-Inhalte überspringen
            skip = ["musterlösung", "trainer-hinweis", "trainer-tipp", "trainertipp",
                    "evaluation", "feedbackbogen", "beobachtungsbogen"]
            if any(k in h3_title.lower() for k in skip):
                continue

            # Praxisfall → image-led
            if any(k in h3_title.lower() for k in
                   ["praxisfall", "der fall", "fallstudie", "fallbeispiel"]):
                pf_name2, pf_sit2 = find_praxisfall("### " + h3_title + "\n" + h3_body)
                name  = pf_name2 or h3_title
                sit   = pf_sit2 or clean_md(h3_body[:400])
                slide_10_image_led(prs, name, sit, name, sec_label, page, pt)
                page += 1
                continue

            # Reflexionsfragen → eigene Folie
            if any(k in h3_title.lower() for k in ["reflexion", "fragen"]):
                qs = find_reflexionsfragen("### " + h3_title + "\n" + h3_body)
                if not qs:
                    qs = extract_bullets(h3_body, 5)
                if qs:
                    slide_content(prs, h3_title, qs, sec_label, page, pt)
                    page += 1
                continue

            # Tabelle im Body?
            tbl = parse_md_table(h3_body)
            if tbl:
                ctx = ""
                for para in re.split(r"\n{2,}", h3_body):
                    if not para.startswith("|"):
                        ctx = clean_md(para)[:200]
                        break
                slide_08_table(prs, h3_title, tbl, ctx, sec_label, page, pt)
                page += 1

                # H4 nach der Tabelle
                h4s = split_sections(h3_body, 4)
                for h4 in h4s:
                    if any(k in h4["title"].lower() for k in skip):
                        continue
                    bullets = extract_bullets(h4["body"])
                    if bullets:
                        slide_content(prs, clean_md(h4["title"]), bullets,
                                      sec_label, page, pt)
                        page += 1
                continue

            # H4-Unterabschnitte?
            h4s = split_sections(h3_body, 4)
            if h4s:
                # Headline-Folie für H3
                intro = re.split(r"^####", h3_body, maxsplit=1, flags=re.MULTILINE)[0]
                intro_bullets = extract_bullets(intro, 3)
                left_text  = " ".join(intro_bullets) if intro_bullets else clean_md(intro[:300])
                slide_04_headline(prs, h3_title, left_text, "", sec_label, page, pt)
                page += 1

                for h4 in h4s:
                    if any(k in h4["title"].lower() for k in skip):
                        continue
                    h4_title  = clean_md(h4["title"])
                    h4_body   = h4["body"]

                    tbl4 = parse_md_table(h4_body)
                    if tbl4:
                        slide_08_table(prs, h4_title, tbl4, "", sec_label, page, pt)
                    else:
                        bullets = extract_bullets(h4_body, 6)
                        if bullets:
                            slide_content(prs, h4_title, bullets, sec_label, page, pt)
                    page += 1
                continue

            # Keine Unterstruktur: Bullets
            bullets = extract_bullets(h3_body, 6)
            if bullets:
                # Viele Bullets → zwei Spalten
                if len(bullets) >= 4:
                    half = len(bullets) // 2
                    slide_05_two_col(
                        prs, h3_title,
                        [""] + bullets[:half],
                        [""] + bullets[half:],
                        sec_label, page, pt)
                else:
                    slide_content(prs, h3_title, bullets, sec_label, page, pt)
                page += 1

    # CLOSING
    slide_12_closing(prs, m, page, str(page))

    out = OUTPUT_DIR / f"{module_id.upper()}.pptx"
    prs.save(out)
    total = len(prs.slides)
    print(f"✓  {out}  ({total} Folien)")
    return out


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    if len(sys.argv) < 2:
        print("Verwendung: python3 generate_praesentation.py M01|all")
        sys.exit(1)

    arg = sys.argv[1].strip().lower()
    if arg == "all":
        ids = sorted(p.stem for p in MODULES_DIR.glob("*.md"))
        for mid in ids:
            try:
                build_prs(mid)
            except Exception as e:
                print(f"✗  {mid}: {e}")
    else:
        build_prs(arg)

if __name__ == "__main__":
    main()
